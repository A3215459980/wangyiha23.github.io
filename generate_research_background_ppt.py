from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


COLORS = {
    "navy": rgb("#183153"),
    "blue": rgb("#2563EB"),
    "light_blue": rgb("#EAF2FF"),
    "teal": rgb("#0F766E"),
    "light_teal": rgb("#E7F8F4"),
    "orange": rgb("#EA580C"),
    "light_orange": rgb("#FFF2E8"),
    "red": rgb("#DC2626"),
    "light_red": rgb("#FDECEC"),
    "purple": rgb("#6D28D9"),
    "light_purple": rgb("#F3EEFF"),
    "green": rgb("#15803D"),
    "light_green": rgb("#EAF8EE"),
    "gray": rgb("#475569"),
    "light_gray": rgb("#F8FAFC"),
    "mid_gray": rgb("#CBD5E1"),
    "dark": rgb("#0F172A"),
    "white": rgb("#FFFFFF"),
}


def set_text_style(run, size, bold=False, color=None, font_name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=None, align=PP_ALIGN.LEFT, fill=None, line=None,
                radius=False, margin=0.08, font_name="Microsoft YaHei"):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    box = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill if fill is not None else COLORS["white"]
    box.line.color.rgb = line if line is not None else box.fill.fore_color.rgb
    if line is None and fill is None:
        box.line.fill.background()
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin / 1.5)
    text_frame.margin_bottom = Inches(margin / 1.5)
    for idx, paragraph_text in enumerate(text.split("\n")):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = paragraph_text
        set_text_style(
            run,
            size=size,
            bold=bold if idx == 0 else False,
            color=color or COLORS["dark"],
            font_name=font_name,
        )
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
    return box


def add_bullet_box(slide, left, top, width, height, title, bullets, accent, fill):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_text_style(r, 19, bold=True, color=accent)

    for bullet in bullets:
        paragraph = tf.add_paragraph()
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(2)
        paragraph.space_after = Pt(1)
        run = paragraph.add_run()
        run.text = bullet
        set_text_style(run, 14.5, color=COLORS["dark"])
    return box


def add_connector(slide, x1, y1, x2, y2, color):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def add_header(slide, title, subtitle):
    add_textbox(
        slide,
        0,
        0,
        SLIDE_W,
        0.55,
        "",
        fill=COLORS["navy"],
        line=COLORS["navy"],
    )
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.15), Inches(8.5), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_text_style(r, 26, bold=True, color=COLORS["white"])

    sub_box = slide.shapes.add_textbox(Inches(9.0), Inches(0.18), Inches(3.8), Inches(0.28))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = subtitle
    set_text_style(r2, 10.5, color=COLORS["white"])


def build_slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "研究背景：大规模 MU-MIMO 检测的性能-复杂度瓶颈",
        "组会汇报 | Background",
    )

    add_textbox(
        slide,
        0.45,
        0.78,
        12.4,
        0.65,
        "核心问题：随着用户数增加，MUI 与病态信道共同放大检测难度，系统亟需在“低复杂度”与“高可靠检测”之间取得可实现的平衡。",
        size=18,
        bold=True,
        color=COLORS["navy"],
        fill=COLORS["light_blue"],
        line=COLORS["blue"],
        radius=True,
        margin=0.12,
    )

    add_bullet_box(
        slide,
        0.45,
        1.65,
        4.0,
        3.55,
        "1) 场景挑战：系统规模上升后检测更难",
        [
            "5G/6G 追求更高频谱效率与容量，MU-MIMO 成为关键支撑技术。",
            "基站服务用户数增加后，多用户干扰（MUI）显著增强。",
            "信道矩阵病态或高度相关时，线性分离能力下降，误检更易累积。",
            "接收端检测因此成为制约链路可靠性的核心瓶颈。",
        ],
        COLORS["blue"],
        COLORS["light_blue"],
    )

    add_bullet_box(
        slide,
        4.67,
        1.65,
        4.02,
        3.55,
        "2) 传统检测方法：两端都不理想",
        [
            "MMSE 等线性检测复杂度较低，工程实现友好。",
            "但在强干扰、高相关信道下性能明显退化。",
            "MAP 检测理论上最优，可最小化符号错误率。",
            "其复杂度随维度指数增长，难以直接用于大规模系统。",
        ],
        COLORS["orange"],
        COLORS["light_orange"],
    )

    add_bullet_box(
        slide,
        8.89,
        1.65,
        3.98,
        3.55,
        "3) EP 折中方案：有效但仍有上限",
        [
            "EP 通过矩匹配将复杂后验近似为高斯分布，并迭代求解。",
            "它在性能与复杂度之间提供了有吸引力的折中。",
            "但“独立高斯近似”忽略变量高阶相关性，强干扰场景下信息损失明显。",
            "高维条件下还存在数值稳定性与收敛性挑战。",
        ],
        COLORS["teal"],
        COLORS["light_teal"],
    )

    add_textbox(
        slide,
        0.65,
        5.45,
        12.0,
        1.18,
        "科学问题凝练\n如何突破 EP 中独立高斯近似带来的性能上限，并在大规模 MU-MIMO 条件下继续保持近线性、可部署的检测复杂度？",
        size=18,
        bold=True,
        color=COLORS["red"],
        fill=COLORS["light_red"],
        line=COLORS["red"],
        radius=True,
        margin=0.16,
    )
    return slide


def build_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "研究进展、方法局限与本文切入点",
        "组会汇报 | Motivation",
    )

    add_textbox(
        slide,
        0.45,
        0.78,
        6.15,
        0.45,
        "已有研究主线：从模型驱动深度展开，到 GNN 辅助概率推断",
        size=17,
        bold=True,
        color=COLORS["navy"],
        fill=COLORS["light_gray"],
        line=COLORS["mid_gray"],
        radius=True,
        margin=0.10,
    )

    steps = [
        (
            0.55,
            "深度展开",
            "OAMP-Net2：\n将迭代检测映射为多层网络，引入可学习参数，兼顾机理可解释性与检测精度。",
            COLORS["blue"],
            COLORS["light_blue"],
        ),
        (
            2.65,
            "GNN + EP",
            "GEPNet / GCEPNet：\n把图结构嵌入 EP 或腔分布估计，用图特征缓解独立高斯近似偏差。",
            COLORS["teal"],
            COLORS["light_teal"],
        ),
        (
            4.75,
            "图辅助 AMP",
            "GNN-Assisted BiG-AMP：\n在消息传递循环中补偿边际似然近似误差，提升鲁棒性。",
            COLORS["purple"],
            COLORS["light_purple"],
        ),
    ]

    for idx, (left, title, body, accent, fill) in enumerate(steps):
        add_textbox(
            slide,
            left,
            1.45,
            1.5,
            0.5,
            title,
            size=15.5,
            bold=True,
            color=COLORS["white"],
            fill=accent,
            line=accent,
            radius=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left - 0.05,
            2.0,
            1.6,
            1.45,
            body,
            size=12.5,
            color=COLORS["dark"],
            fill=fill,
            line=accent,
            radius=True,
            margin=0.09,
        )
        if idx < 2:
            add_connector(slide, left + 1.62, 2.72, left + 1.98, 2.72, COLORS["mid_gray"])

    add_textbox(
        slide,
        0.55,
        3.78,
        6.0,
        1.85,
        "现有 GNN 方法的主要局限\n• 图卷积依赖局部邻域聚合，难以建模大规模 MIMO 中的长程空间依赖。\n• 网络加深后易出现过平滑，节点特征区分度下降。\n• 在大规模阵列下，图运算开销仍然偏高，限制工程部署。",
        size=15,
        bold=True,
        color=COLORS["orange"],
        fill=COLORS["light_orange"],
        line=COLORS["orange"],
        radius=True,
        margin=0.14,
    )

    add_textbox(
        slide,
        6.85,
        0.78,
        5.95,
        0.45,
        "本文切入点：引入 DIFFormer，增强全局依赖建模能力",
        size=17,
        bold=True,
        color=COLORS["navy"],
        fill=COLORS["light_green"],
        line=COLORS["green"],
        radius=True,
        margin=0.10,
    )

    add_textbox(
        slide,
        6.9,
        1.45,
        5.85,
        2.2,
        "DIFFormer 的关键优势\n1. 通过最小化表征全局一致性的能量函数，构造扩散诱导注意力。\n2. 不受限于固定局部邻域，可在任意节点对之间动态传播信息。\n3. 以近线性复杂度聚合全局干扰特征，兼顾表达力与可扩展性。\n4. 通过解耦局部结构与全局交互，缓解传统深层 GNN 的过平滑问题。",
        size=14.5,
        bold=True,
        color=COLORS["green"],
        fill=COLORS["light_green"],
        line=COLORS["green"],
        radius=True,
        margin=0.14,
    )

    add_textbox(
        slide,
        6.9,
        3.95,
        5.85,
        1.68,
        "拟议框架：D-EPNet\n将 DIFFormer 深度嵌入 EP 的迭代推断骨架，用全局扩散能力动态修正腔分布估计偏差，从数学本质上突破独立高斯近似的限制，为 6G 高性能检测提供新方案。",
        size=15.5,
        bold=True,
        color=COLORS["red"],
        fill=COLORS["light_red"],
        line=COLORS["red"],
        radius=True,
        margin=0.14,
    )

    refs = (
        "参考文献索引：[1] GEPNet  [2] GCEPNet  [3] He 等模型驱动深度学习  "
        "[4] Scotti 等 GNN-MRF 检测  [5] GNN-Assisted BiG-AMP  [6] DIFFormer"
    )
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.86), Inches(12.1), Inches(0.25))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = refs
    set_text_style(r, 9.5, color=COLORS["gray"])
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_slide_one(prs)
    build_slide_two(prs)

    output_path = "mu_mimo_research_background_group_meeting.pptx"
    prs.save(output_path)
    print(f"Saved PPT to {output_path}")


if __name__ == "__main__":
    main()
